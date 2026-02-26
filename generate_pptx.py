#!/usr/bin/env python3
"""Generate PowerPoint presentation for AI News Collector Bot."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ── Colour palette ──────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT_BLUE    = RGBColor(0x1A, 0x73, 0xE8)
ACCENT_GREEN   = RGBColor(0x0D, 0x96, 0x52)
ACCENT_ORANGE  = RGBColor(0xE8, 0x71, 0x0A)
ACCENT_PURPLE  = RGBColor(0x7B, 0x1F, 0xA2)
ACCENT_RED     = RGBColor(0xD9, 0x30, 0x25)
ACCENT_TEAL    = RGBColor(0x00, 0x96, 0x88)
LIGHT_BLUE     = RGBColor(0xE8, 0xF0, 0xFE)
LIGHT_GREEN    = RGBColor(0xE6, 0xF4, 0xEA)
LIGHT_ORANGE   = RGBColor(0xFE, 0xF7, 0xE0)
LIGHT_PURPLE   = RGBColor(0xF3, 0xE8, 0xFD)
LIGHT_RED      = RGBColor(0xFC, 0xE8, 0xE6)
LIGHT_TEAL     = RGBColor(0xE0, 0xF2, 0xF1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ────────────────────────────────────────────

def set_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_page_number(slide, num):
    """Add page number in bottom-right."""
    txBox = slide.shapes.add_textbox(Inches(12.1), Inches(7.0), Inches(1.0), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title_text(slide, text, left, top, width, height, font_size=36,
                   bold=True, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, font_size=18,
                  color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=DARK_GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tf


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=True, text_align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = text_align
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top, color=MID_GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.end_x = end_left
    connector.end_y = end_top
    return connector


def add_icon_card(slide, left, top, width, height, icon_text, title, desc,
                  accent_color, bg_color):
    """Add a card with icon circle, title, and description."""
    # Background card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Icon circle
    circle_size = Inches(0.7)
    cx = left + Inches(0.2)
    cy = top + Inches(0.2)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, circle_size, circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.paragraphs[0].text = icon_text
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_title_text(slide, title,
                   left + Inches(1.1), top + Inches(0.15),
                   width - Inches(1.3), Inches(0.5),
                   font_size=18, bold=True, color=DARK_GRAY)
    # Description
    add_body_text(slide, desc,
                  left + Inches(0.2), top + Inches(0.9),
                  width - Inches(0.4), height - Inches(1.0),
                  font_size=13, color=MID_GRAY)


def add_flow_box(slide, left, top, width, height, text, fill_color, font_color=WHITE, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = font_color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_simple_arrow_right(slide, x, y, length=Inches(0.4), color=MID_GRAY):
    """Draw a right-pointing arrow using a connector."""
    connector = slide.shapes.add_connector(1, x, y, x + length, y)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def add_simple_arrow_down(slide, x, y, length=Inches(0.4), color=MID_GRAY):
    connector = slide.shapes.add_connector(1, x, y, x, y + length)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def add_section_header(slide, text, subtitle=""):
    """Add a blue accent bar + title for section slides."""
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(0.5), Inches(0.12), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_title_text(slide, text,
                   Inches(1.0), Inches(0.45), Inches(11), Inches(0.9),
                   font_size=36, bold=True, color=DARK_GRAY)
    if subtitle:
        add_body_text(slide, subtitle,
                      Inches(1.0), Inches(1.3), Inches(11), Inches(0.5),
                      font_size=18, color=MID_GRAY)


# ================================================================
# SLIDE 1: Title Slide
# ================================================================
page = 1
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_white_bg(slide)

# Top accent line
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(0), SLIDE_W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Main title
add_title_text(slide, "AI News Collector Bot",
               Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
               font_size=52, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Subtitle
add_body_text(slide, "AIニュースを自動収集し、ストーリー形式で毎朝届ける",
              Inches(1.5), Inches(3.1), Inches(10), Inches(0.7),
              font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# 4 feature icons at bottom
icons = [
    ("📡", "自動収集", ACCENT_BLUE),
    ("📖", "ストーリー変換", ACCENT_GREEN),
    ("📧", "マルチ配信", ACCENT_ORANGE),
    ("💡", "ナレッジ蓄積", ACCENT_PURPLE),
]
start_x = Inches(2.5)
for i, (emoji, label, color) in enumerate(icons):
    cx = start_x + Inches(i * 2.3)
    cy = Inches(4.8)
    # Circle
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.9), Inches(0.9))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.paragraphs[0].text = emoji
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Label
    add_body_text(slide, label, cx - Inches(0.3), cy + Inches(1.05),
                  Inches(1.5), Inches(0.4),
                  font_size=13, color=MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, page)


# ================================================================
# SLIDE 2: 想定課題
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "想定課題", "AI情報のキャッチアップに、こんな悩みはありませんか？")

problems = [
    ("😰", "情報過多", "AIニュースは毎日大量に発信されるが、\n全てに目を通す時間がない", ACCENT_RED, LIGHT_RED),
    ("🔍", "ソース分散", "RSS・ニュースサイト・SNSなど情報源が\nバラバラで一元管理できない", ACCENT_ORANGE, LIGHT_ORANGE),
    ("📋", "記録が残らない", "読んだニュースが流れてしまい、\n後から検索・振り返りができない", ACCENT_PURPLE, LIGHT_PURPLE),
    ("😴", "読みにくい", "技術記事は専門的で読みにくく、\n要点の把握に時間がかかる", ACCENT_TEAL, LIGHT_TEAL),
]

for i, (icon, title, desc, accent, bg) in enumerate(problems):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + Inches(col * 6.2)
    top = Inches(2.2) + Inches(row * 2.5)
    w = Inches(5.8)
    h = Inches(2.2)
    add_icon_card(slide, left, top, w, h, icon, title, desc, accent, bg)

add_page_number(slide, page)


# ================================================================
# SLIDE 3: ソリューション概要
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "AI News Collector Bot とは",
                   "毎朝6時に自動実行。AIニュース3本を厳選し、読みやすいストーリー形式で配信します。")

# Value proposition cards
values = [
    ("📡", "自動収集", "RSS / Webスクレイピング /\nNewsAPIから毎日3本を厳選",
     ACCENT_BLUE, LIGHT_BLUE),
    ("📖", "ストーリー変換", "STAR・英雄の旅 等の\nフレームワークで読みやすく変換",
     ACCENT_GREEN, LIGHT_GREEN),
    ("📧", "マルチ配信", "Gmail (HTML) と\nLINE Notify で毎朝届く",
     ACCENT_ORANGE, LIGHT_ORANGE),
    ("💡", "ナレッジ蓄積", "Markdown保存・タグ検索\nリアクション・月次サマリ",
     ACCENT_PURPLE, LIGHT_PURPLE),
]

for i, (icon, title, desc, accent, bg) in enumerate(values):
    left = Inches(0.6) + Inches(i * 3.15)
    top = Inches(2.5)
    w = Inches(2.9)
    h = Inches(3.8)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Top accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        left, top, w, Inches(0.06))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    # Icon
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    left + Inches(0.9), top + Inches(0.4),
                                    Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_title_text(slide, title,
                   left + Inches(0.2), top + Inches(1.6), w - Inches(0.4), Inches(0.5),
                   font_size=20, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    # Desc
    add_body_text(slide, desc,
                  left + Inches(0.2), top + Inches(2.2), w - Inches(0.4), Inches(1.2),
                  font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow connectors between cards
for i in range(3):
    ax = Inches(0.6) + Inches((i + 1) * 3.15) - Inches(0.15)
    ay = Inches(4.4)
    add_simple_arrow_right(slide, ax, ay, Inches(0.15), ACCENT_BLUE)

add_page_number(slide, page)


# ================================================================
# SLIDE 4: システム構成図
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "システム構成図", "5つのモジュールが連携して動作するパイプラインアーキテクチャ")

# Module boxes
modules = [
    ("Collector\n(ニュース収集)", Inches(0.5), Inches(2.8), Inches(2.2), Inches(1.4), ACCENT_BLUE),
    ("Selector\n(3本厳選)", Inches(3.2), Inches(2.8), Inches(2.0), Inches(1.4), ACCENT_TEAL),
    ("Writer\n(ストーリー変換)", Inches(5.7), Inches(2.8), Inches(2.2), Inches(1.4), ACCENT_GREEN),
    ("Delivery\n(Gmail / LINE)", Inches(8.4), Inches(2.8), Inches(2.2), Inches(1.4), ACCENT_ORANGE),
    ("Knowledge Base\n(Markdown保存)", Inches(10.9), Inches(2.8), Inches(2.0), Inches(1.4), ACCENT_PURPLE),
]

for text, l, t, w, h, color in modules:
    add_flow_box(slide, l, t, w, h, text, color, WHITE, 14)

# Arrows between modules
arrow_positions = [
    (Inches(2.7), Inches(3.5), Inches(3.2), Inches(3.5)),
    (Inches(5.2), Inches(3.5), Inches(5.7), Inches(3.5)),
    (Inches(7.9), Inches(3.5), Inches(8.4), Inches(3.5)),
    (Inches(10.6), Inches(3.5), Inches(10.9), Inches(3.5)),
]
for sx, sy, ex, ey in arrow_positions:
    c = slide.shapes.add_connector(1, sx, sy, ex, ey)
    c.line.color.rgb = MID_GRAY
    c.line.width = Pt(2.5)

# Sub-labels under modules
sub_labels = [
    ("RSS / Web /\nNewsAPI", Inches(0.5), Inches(4.4), Inches(2.2)),
    ("重複排除\nスコアリング", Inches(3.2), Inches(4.4), Inches(2.0)),
    ("STAR / 英雄の旅\nBefore-After / PAS", Inches(5.7), Inches(4.4), Inches(2.2)),
    ("HTML変換\nリアクション付き", Inches(8.4), Inches(4.4), Inches(2.2)),
    ("YAML Frontmatter\nタグ / 検索", Inches(10.9), Inches(4.4), Inches(2.0)),
]
for text, l, t, w in sub_labels:
    add_body_text(slide, text, l, t, w, Inches(0.8),
                  font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Feedback Server box at bottom
add_flow_box(slide, Inches(4.5), Inches(5.6), Inches(4.5), Inches(1.2),
             "Feedback Server (FastAPI / Port 8321)\nリアクション受付 ← → ナレッジベース更新",
             ACCENT_RED, WHITE, 13)

# Arrow from feedback to knowledge base
c = slide.shapes.add_connector(1, Inches(9.0), Inches(6.2), Inches(11.9), Inches(4.2))
c.line.color.rgb = MID_GRAY
c.line.width = Pt(1.5)

# Cron label at top
add_flow_box(slide, Inches(4.8), Inches(1.7), Inches(3.5), Inches(0.7),
             "⏰  Cron 毎朝 6:00 JST 自動実行", LIGHT_GRAY, DARK_GRAY, 13)
c = slide.shapes.add_connector(1, Inches(6.55), Inches(2.4), Inches(6.55), Inches(2.8))
c.line.color.rgb = MID_GRAY
c.line.width = Pt(1.5)

add_page_number(slide, page)


# ================================================================
# SLIDE 5: 主要機能一覧
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "主要機能一覧", "4つの柱で構成される統合AIニュースシステム")

features = [
    ("1", "ニュース自動収集", [
        "RSS (feedparser) による定期フィード取得",
        "BeautifulSoup による Webスクレイピング",
        "NewsAPI 連携（キーワード: AI, LLM, GPT 等）",
        "重複排除 & 新鮮度スコアリングで3本厳選",
    ], ACCENT_BLUE, LIGHT_BLUE),
    ("2", "ストーリーテリング変換", [
        "Claude API でニュースを物語形式に変換",
        "4つのフレームワーク: STAR / 英雄の旅 / B-A-B / PAS",
        "記事内容に応じて最適フレームワーク自動選定",
        "Jinja2テンプレートでMarkdown生成",
    ], ACCENT_GREEN, LIGHT_GREEN),
    ("3", "マルチチャネル配信", [
        "Gmail API (OAuth2) による HTML メール送信",
        "LINE Notify によるテキスト通知",
        "リアクションボタン付きメール (⭐👍🤔📌)",
        "HTMLテンプレートで見やすいデザイン",
    ], ACCENT_ORANGE, LIGHT_ORANGE),
    ("4", "ナレッジベース & フィードバック", [
        "YAML Frontmatter付きMarkdownで知識蓄積",
        "FastAPI サーバー (Port 8321) でリアクション受付",
        "タグ検索・全文検索・月次サマリ自動生成",
        "カウンターベースの4段階リアクションシステム",
    ], ACCENT_PURPLE, LIGHT_PURPLE),
]

for i, (num, title, items, accent, bg) in enumerate(features):
    left = Inches(0.5)
    top = Inches(2.2) + Inches(i * 1.25)
    w = Inches(12.3)
    h = Inches(1.1)

    # Row background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.fill.background()
    card.shadow.inherit = False

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    left + Inches(0.15), top + Inches(0.18),
                                    Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_title_text(slide, title,
                   left + Inches(1.1), top + Inches(0.05),
                   Inches(3.0), Inches(0.45),
                   font_size=18, bold=True, color=accent)

    # Feature bullets (as a single line)
    bullet_text = "  |  ".join(items)
    add_body_text(slide, bullet_text,
                  left + Inches(1.1), top + Inches(0.5),
                  Inches(11.0), Inches(0.55),
                  font_size=12, color=MID_GRAY)

add_page_number(slide, page)


# ================================================================
# SLIDE 6: ニュース自動収集の詳細
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "① ニュース自動収集", "3つのソースからAIニュースを収集し、最も価値の高い3本を自動選定")

# Three source cards
sources = [
    ("RSS\nフィード", "feedparserライブラリで\n主要AIメディアのRSSを\n定期取得\n\n・TechCrunch AI\n・MIT Tech Review\n・AI専門ブログ 等",
     ACCENT_BLUE, LIGHT_BLUE),
    ("Web\nスクレイピング", "BeautifulSoupで\nWebページから記事を\n自動抽出\n\n・タイトル / 本文\n・公開日 / ソース名",
     ACCENT_TEAL, LIGHT_TEAL),
    ("NewsAPI\n連携", "NewsAPI.orgから\nキーワード検索で取得\n\n・AI / LLM / GPT\n・機械学習 / 深層学習\n・最新24時間の記事",
     ACCENT_GREEN, LIGHT_GREEN),
]

for i, (title, desc, accent, bg) in enumerate(sources):
    left = Inches(0.6) + Inches(i * 3.6)
    top = Inches(2.3)
    w = Inches(3.2)
    h = Inches(3.8)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)
    card.shadow.inherit = False

    add_rounded_rect(slide, left + Inches(0.6), top + Inches(0.2),
                     Inches(2.0), Inches(0.9), accent, title, 16, WHITE, True)

    add_body_text(slide, desc,
                  left + Inches(0.2), top + Inches(1.3),
                  w - Inches(0.4), Inches(2.3),
                  font_size=13, color=MID_GRAY)

# Selection process box on right
sel_left = Inches(11.2)
add_flow_box(slide, sel_left, Inches(2.3), Inches(1.8), Inches(0.8),
             "重複排除", ACCENT_RED, WHITE, 13)
add_simple_arrow_down(slide, sel_left + Inches(0.9), Inches(3.1), Inches(0.3), MID_GRAY)
add_flow_box(slide, sel_left, Inches(3.4), Inches(1.8), Inches(0.8),
             "スコアリング", ACCENT_ORANGE, WHITE, 13)
add_simple_arrow_down(slide, sel_left + Inches(0.9), Inches(4.2), Inches(0.3), MID_GRAY)
add_flow_box(slide, sel_left, Inches(4.5), Inches(1.8), Inches(0.8),
             "Top 3 選定", ACCENT_GREEN, WHITE, 13)

# Arrow from sources to selection
c = slide.shapes.add_connector(1, Inches(10.6), Inches(4.0), Inches(11.2), Inches(4.0))
c.line.color.rgb = MID_GRAY
c.line.width = Pt(2)

add_page_number(slide, page)


# ================================================================
# SLIDE 7: ストーリーテリング変換
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "② ストーリーテリング変換",
                   "Claude APIが記事を分析し、最適なフレームワークで物語形式に変換")

frameworks = [
    ("STAR法", "Situation → Task → Action → Result",
     "企業のAI導入事例など\n成果を伝えたい記事に最適",
     ACCENT_BLUE, LIGHT_BLUE),
    ("英雄の旅", "日常 → 冒険 → 試練 → 帰還",
     "技術革新・ブレイクスルー系の\n記事に最適",
     ACCENT_GREEN, LIGHT_GREEN),
    ("Before/After\n/Bridge", "課題 → 理想 → 解決策",
     "業務改善・効率化系の\n記事に最適",
     ACCENT_ORANGE, LIGHT_ORANGE),
    ("PAS法", "Problem → Agitate → Solution",
     "問題解決・課題提起系の\n記事に最適",
     ACCENT_PURPLE, LIGHT_PURPLE),
]

for i, (name, flow, desc, accent, bg) in enumerate(frameworks):
    left = Inches(0.5) + Inches(i * 3.2)
    top = Inches(2.5)
    w = Inches(2.95)
    h = Inches(4.2)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Name
    add_rounded_rect(slide, left + Inches(0.3), top + Inches(0.2),
                     Inches(2.35), Inches(0.7), accent, name, 18, WHITE, True)

    # Flow
    add_body_text(slide, flow,
                  left + Inches(0.15), top + Inches(1.1),
                  w - Inches(0.3), Inches(0.8),
                  font_size=12, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    # Separator line
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 left + Inches(0.5), top + Inches(2.0),
                                 Inches(1.95), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = LIGHT_GRAY
    sep.line.fill.background()

    # Description
    add_body_text(slide, desc,
                  left + Inches(0.15), top + Inches(2.2),
                  w - Inches(0.3), Inches(1.5),
                  font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note
add_body_text(slide, "💡 記事の内容をキーワード分析 → Claude APIが最適フレームワークを自動選定 → Jinja2テンプレートでMarkdown生成",
              Inches(1.0), Inches(6.9), Inches(11), Inches(0.4),
              font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, page)


# ================================================================
# SLIDE 8: マルチチャネル配信
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "③ マルチチャネル配信",
                   "Gmail (HTML) と LINE Notify で、毎朝読みやすいニュースダイジェストを届ける")

# Gmail section
add_rounded_rect(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.5),
                 LIGHT_BLUE, "", 14, DARK_GRAY)

add_title_text(slide, "📧  Gmail (HTML メール)",
               Inches(0.9), Inches(2.5), Inches(5.2), Inches(0.5),
               font_size=22, bold=True, color=ACCENT_BLUE)

gmail_features = [
    "✅  OAuth2 認証による安全な送信",
    "✅  HTMLテンプレートで美しいデザイン",
    "✅  リアクションボタン埋め込み",
    "✅  ⭐ 素晴らしい / 👍 良い / 🤔 微妙 / 📌 後で読む",
    "✅  ワンクリックでフィードバック送信",
    "✅  レスポンシブ対応メール",
]
add_bullet_list(slide, gmail_features,
                Inches(1.2), Inches(3.2), Inches(4.8), Inches(3.2),
                font_size=14, color=DARK_GRAY)

# LINE section
add_rounded_rect(slide, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.5),
                 LIGHT_GREEN, "", 14, DARK_GRAY)

add_title_text(slide, "💬  LINE Notify",
               Inches(7.2), Inches(2.5), Inches(5.2), Inches(0.5),
               font_size=22, bold=True, color=ACCENT_GREEN)

line_features = [
    "✅  LINE Notify API でプッシュ通知",
    "✅  テキスト形式で要点を簡潔に配信",
    "✅  各記事のタイトル + 要約 + リンク",
    "✅  スマホですぐに確認可能",
    "✅  通勤中・移動中に最適",
    "✅  設定ONで個人LINE / グループに送信",
]
add_bullet_list(slide, line_features,
                Inches(7.5), Inches(3.2), Inches(4.8), Inches(3.2),
                font_size=14, color=DARK_GRAY)

add_page_number(slide, page)


# ================================================================
# SLIDE 9: フィードバック & ナレッジベース
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "④ フィードバック & ナレッジベース",
                   "リアクションで評価 → 知識を蓄積 → 検索・振り返りが可能に")

# Left: Reaction system
add_rounded_rect(slide, Inches(0.6), Inches(2.3), Inches(6.0), Inches(4.8),
                 LIGHT_ORANGE, "", 14, DARK_GRAY)

add_title_text(slide, "リアクションシステム",
               Inches(0.9), Inches(2.5), Inches(5.4), Inches(0.5),
               font_size=20, bold=True, color=ACCENT_ORANGE)

reactions = [
    ("⭐", "素晴らしい (excellent)", "カウンター +1"),
    ("👍", "良い (good)", "カウンター +1"),
    ("🤔", "微妙 (so_so)", "カウンター +1"),
    ("📌", "後で読む (read_later)", "カウンター +1"),
]

for i, (emoji, label, counter) in enumerate(reactions):
    y = Inches(3.2) + Inches(i * 0.6)
    # Emoji circle
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.0), y, Inches(0.5), Inches(0.5))
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT_ORANGE
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.paragraphs[0].text = emoji
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_body_text(slide, f"{label}    →    {counter}",
                  Inches(1.7), y, Inches(4.5), Inches(0.5),
                  font_size=14, color=DARK_GRAY)

add_body_text(slide, "RESTful API:  GET /api/reaction/{date}/{story_id}/{type}",
              Inches(1.0), Inches(5.7), Inches(5.2), Inches(0.4),
              font_size=12, color=ACCENT_ORANGE, bold=True)

add_body_text(slide, "FastAPI サーバー  |  Port 8321  |  自動Frontmatter更新",
              Inches(1.0), Inches(6.2), Inches(5.2), Inches(0.4),
              font_size=12, color=MID_GRAY)

# Right: Knowledge base
add_rounded_rect(slide, Inches(7.0), Inches(2.3), Inches(6.0), Inches(4.8),
                 LIGHT_PURPLE, "", 14, DARK_GRAY)

add_title_text(slide, "ナレッジベース",
               Inches(7.3), Inches(2.5), Inches(5.4), Inches(0.5),
               font_size=20, bold=True, color=ACCENT_PURPLE)

kb_features = [
    "📁  YAML Frontmatter付き Markdown で保存",
    "🏷️  自動タグ付け & カテゴリ分類",
    "🔍  タグ検索 + 全文検索に対応",
    "📊  月次サマリ自動生成",
    "📅  日付ベースのディレクトリ構造",
    "🔒  ファイルロックによる安全な並行更新",
]
add_bullet_list(slide, kb_features,
                Inches(7.5), Inches(3.2), Inches(5.0), Inches(3.5),
                font_size=14, color=DARK_GRAY)

add_page_number(slide, page)


# ================================================================
# SLIDE 10: 日次処理フロー
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "日次処理フロー", "毎朝6:00 JSTに自動実行される一連の処理パイプライン")

# Flow steps
steps = [
    ("6:00\nCron起動", ACCENT_RED),
    ("RSS/Web/API\nニュース収集", ACCENT_BLUE),
    ("重複排除\n3本厳選", ACCENT_TEAL),
    ("フレームワーク\n自動選定", ACCENT_GREEN),
    ("Claude API\nストーリー変換", ACCENT_GREEN),
    ("Markdown\n生成・保存", ACCENT_PURPLE),
    ("Gmail\nHTML送信", ACCENT_ORANGE),
    ("LINE\nNotify送信", ACCENT_ORANGE),
]

box_w = Inches(1.35)
box_h = Inches(1.1)
start_x = Inches(0.4)
y_row1 = Inches(2.8)
y_row2 = Inches(5.0)

for i, (text, color) in enumerate(steps):
    if i < 4:
        x = start_x + Inches(i * 3.2)
        y = y_row1
    else:
        x = start_x + Inches((7 - i) * 3.2)
        y = y_row2

    add_flow_box(slide, x, y, box_w, box_h, text, color, WHITE, 12)

    # Step number
    add_body_text(slide, f"Step {i + 1}",
                  x, y - Inches(0.3), box_w, Inches(0.3),
                  font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# Horizontal arrows row 1
for i in range(3):
    ax = start_x + box_w + Inches(i * 3.2)
    ay = y_row1 + box_h / 2
    c = slide.shapes.add_connector(1, ax, ay, ax + Inches(3.2) - box_w, ay)
    c.line.color.rgb = MID_GRAY
    c.line.width = Pt(2)

# Down arrow from step 4 to step 5
c = slide.shapes.add_connector(1,
                               start_x + Inches(3 * 3.2) + box_w / 2, y_row1 + box_h,
                               start_x + Inches(3 * 3.2) + box_w / 2, y_row2)
c.line.color.rgb = MID_GRAY
c.line.width = Pt(2)

# Horizontal arrows row 2 (right to left)
for i in range(3):
    ax = start_x + Inches((7 - 4 - i) * 3.2) + box_w
    ay = y_row2 + box_h / 2
    c = slide.shapes.add_connector(1,
                                   ax, ay,
                                   ax + Inches(3.2) - box_w, ay)
    c.line.color.rgb = MID_GRAY
    c.line.width = Pt(2)

# Arrow labels
add_body_text(slide, "▶ ▶ ▶  ▶ ▶ ▶  ▶ ▶ ▶",
              Inches(2.0), Inches(2.4), Inches(9), Inches(0.3),
              font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, page)


# ================================================================
# SLIDE 11: 活用シーン
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "活用シーン", "様々なビジネスシーンで AI 情報のキャッチアップを効率化")

scenes = [
    ("👨‍💻", "エンジニア個人", "毎朝の情報収集を自動化\n通勤中にLINEでAI動向をチェック\nナレッジとして蓄積・検索",
     ACCENT_BLUE, LIGHT_BLUE),
    ("👥", "開発チーム", "チーム全員に同じニュースを共有\nリアクションで興味の傾向を把握\n月次レポートで振り返り",
     ACCENT_GREEN, LIGHT_GREEN),
    ("📈", "マネージャー", "技術トレンドを効率よく把握\nストーリー形式で読みやすい\n意思決定の参考資料に",
     ACCENT_ORANGE, LIGHT_ORANGE),
    ("🎓", "学習者・研究者", "AI分野の最新動向を追跡\nタグ検索で関心領域を深掘り\n学習記録として活用",
     ACCENT_PURPLE, LIGHT_PURPLE),
]

for i, (icon, title, desc, accent, bg) in enumerate(scenes):
    col = i % 2
    row = i // 2
    left = Inches(0.6) + Inches(col * 6.4)
    top = Inches(2.3) + Inches(row * 2.5)
    w = Inches(6.0)
    h = Inches(2.2)
    add_icon_card(slide, left, top, w, h, icon, title, desc, accent, bg)

add_page_number(slide, page)


# ================================================================
# SLIDE 12: ここがスゴい！
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "ここがスゴい！", "他のニュースツールにはない、本ボットの差別化ポイント")

strengths = [
    ("🤖", "AI × ストーリーテリング",
     "単なる要約ではなく、STAR法・英雄の旅など4つのフレームワークで\n「読みたくなる物語」に変換。Claude APIが記事内容を理解し最適な手法を自動選定。",
     ACCENT_BLUE, LIGHT_BLUE),
    ("🔄", "完全自動パイプライン",
     "収集→選定→変換→保存→配信まで全自動。Cron設定一つで毎朝届く。\n--dry-run や --date オプションでテスト・過去日実行も可能。",
     ACCENT_GREEN, LIGHT_GREEN),
    ("💡", "蓄積型ナレッジベース",
     "読み捨てではなく知識が蓄積される。YAML Frontmatter・タグ・全文検索で\nいつでも過去の記事を検索。月次サマリで傾向分析も。",
     ACCENT_PURPLE, LIGHT_PURPLE),
    ("📊", "フィードバックループ",
     "リアクション機能で「読んで終わり」にならない。カウンターベースの評価が\nFrontmatterに記録され、人気記事の傾向把握や今後の改善に活用可能。",
     ACCENT_ORANGE, LIGHT_ORANGE),
]

for i, (icon, title, desc, accent, bg) in enumerate(strengths):
    left = Inches(0.6)
    top = Inches(2.2) + Inches(i * 1.3)
    w = Inches(12.1)
    h = Inches(1.15)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.fill.background()
    card.shadow.inherit = False

    # Icon
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    left + Inches(0.15), top + Inches(0.18),
                                    Inches(0.75), Inches(0.75))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_title_text(slide, title,
                   left + Inches(1.15), top + Inches(0.05),
                   Inches(4.0), Inches(0.4),
                   font_size=18, bold=True, color=accent)

    # Description
    add_body_text(slide, desc,
                  left + Inches(1.15), top + Inches(0.45),
                  Inches(10.5), Inches(0.65),
                  font_size=12, color=MID_GRAY)

add_page_number(slide, page)


# ================================================================
# SLIDE 13: 技術スタック
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "技術スタック", "信頼性の高いオープンソースライブラリとAPI群で構成")

categories = [
    ("言語・ランタイム", ["Python 3.11+"], ACCENT_BLUE),
    ("AI / LLM", ["Claude API (Anthropic)"], ACCENT_GREEN),
    ("ニュース収集", ["feedparser", "BeautifulSoup4", "NewsAPI"], ACCENT_TEAL),
    ("テンプレート", ["Jinja2", "python-frontmatter"], ACCENT_PURPLE),
    ("配信", ["Gmail API (OAuth2)", "LINE Notify API"], ACCENT_ORANGE),
    ("サーバー", ["FastAPI", "uvicorn"], ACCENT_RED),
    ("変換", ["markdown", "Jinja2 HTML"], ACCENT_BLUE),
    ("テスト", ["pytest (219 tests)", "unittest.mock"], ACCENT_GREEN),
]

col_w = Inches(3.0)
row_h = Inches(1.3)

for i, (category, items, color) in enumerate(categories):
    col = i % 4
    row = i // 4
    left = Inches(0.5) + Inches(col * 3.2)
    top = Inches(2.3) + Inches(row * 2.5)

    # Category header
    add_rounded_rect(slide, left, top, col_w, Inches(0.5),
                     color, category, 14, WHITE, True)

    # Items
    items_text = "\n".join([f"• {item}" for item in items])
    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.6),
                                      col_w - Inches(0.2), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

add_page_number(slide, page)


# ================================================================
# SLIDE 14: セットアップ方法
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "導入・セットアップ", "4ステップでセットアップ完了。すぐに使い始められます。")

setup_steps = [
    ("Step 1", "環境構築",
     "git clone → cd ai-news-bot\nbash setup.sh\n(Python仮想環境 + 依存パッケージ自動インストール)",
     ACCENT_BLUE, LIGHT_BLUE),
    ("Step 2", "設定ファイル編集",
     ".env に APIキーを設定:\n• ANTHROPIC_API_KEY\n• NEWS_API_KEY\n• LINE_NOTIFY_TOKEN\n• Gmail OAuth2 credentials",
     ACCENT_GREEN, LIGHT_GREEN),
    ("Step 3", "動作確認",
     "python main.py --dry-run\n(実際の送信なしでテスト実行)\npython main.py --date 2026-02-26\n(特定日のニュースを取得)",
     ACCENT_ORANGE, LIGHT_ORANGE),
    ("Step 4", "自動実行設定",
     "bash scripts/install_cron.sh\n(毎朝6:00 JSTのCronジョブを登録)\npython main.py --server\n(フィードバックサーバー起動)",
     ACCENT_PURPLE, LIGHT_PURPLE),
]

for i, (step, title, desc, accent, bg) in enumerate(setup_steps):
    left = Inches(0.6)
    top = Inches(2.2) + Inches(i * 1.3)
    w = Inches(12.1)
    h = Inches(1.15)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.fill.background()
    card.shadow.inherit = False

    # Step badge
    add_rounded_rect(slide, left + Inches(0.15), top + Inches(0.3),
                     Inches(1.0), Inches(0.5), accent, step, 14, WHITE, True)

    # Title
    add_title_text(slide, title,
                   left + Inches(1.35), top + Inches(0.05),
                   Inches(2.5), Inches(0.4),
                   font_size=18, bold=True, color=accent)

    # Description
    add_body_text(slide, desc.replace("\n", "  |  "),
                  left + Inches(1.35), top + Inches(0.5),
                  Inches(10.5), Inches(0.6),
                  font_size=12, color=MID_GRAY)

add_page_number(slide, page)


# ================================================================
# SLIDE 15: API エンドポイント一覧
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)
add_section_header(slide, "API エンドポイント一覧",
                   "FastAPI サーバー (Port 8321) が提供する RESTful API")

endpoints = [
    ("GET", "/api/reaction/{date}/{story_id}/{type}", "リアクション送信", ACCENT_GREEN),
    ("GET", "/api/stories/{date}", "日付別ストーリー一覧", ACCENT_BLUE),
    ("GET", "/api/stories/{date}/{story_id}", "ストーリー詳細", ACCENT_BLUE),
    ("GET", "/api/search?q=&tag=&min_rating=", "ナレッジベース検索", ACCENT_PURPLE),
    ("GET", "/api/summary/{year}/{month}", "月次サマリ取得", ACCENT_TEAL),
    ("GET", "/api/health", "ヘルスチェック", MID_GRAY),
    ("GET", "/api/stats", "統計情報", MID_GRAY),
]

# Table header
header_y = Inches(2.3)
add_rounded_rect(slide, Inches(0.6), header_y, Inches(1.2), Inches(0.5),
                 DARK_GRAY, "Method", 13, WHITE, True)
add_rounded_rect(slide, Inches(1.8), header_y, Inches(5.5), Inches(0.5),
                 DARK_GRAY, "Endpoint", 13, WHITE, True)
add_rounded_rect(slide, Inches(7.3), header_y, Inches(5.4), Inches(0.5),
                 DARK_GRAY, "説明", 13, WHITE, True)

for i, (method, endpoint, desc, color) in enumerate(endpoints):
    row_y = Inches(2.9) + Inches(i * 0.55)
    bg = WHITE if i % 2 == 0 else LIGHT_BLUE

    # Row background
    row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.6), row_y, Inches(12.1), Inches(0.5))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = bg
    row_bg.line.fill.background()
    row_bg.shadow.inherit = False

    # Method badge
    add_rounded_rect(slide, Inches(0.7), row_y + Inches(0.07),
                     Inches(1.0), Inches(0.36), color, method, 11, WHITE, True)

    # Endpoint
    add_body_text(slide, endpoint,
                  Inches(1.9), row_y + Inches(0.05),
                  Inches(5.3), Inches(0.4),
                  font_size=13, color=DARK_GRAY, bold=True)

    # Description
    add_body_text(slide, desc,
                  Inches(7.4), row_y + Inches(0.05),
                  Inches(5.2), Inches(0.4),
                  font_size=13, color=MID_GRAY)

add_page_number(slide, page)


# ================================================================
# SLIDE 16: まとめ
# ================================================================
page += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_white_bg(slide)

# Top accent line
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(0), SLIDE_W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_title_text(slide, "まとめ",
               Inches(1.0), Inches(0.8), Inches(11), Inches(0.8),
               font_size=40, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Summary points
summary_items = [
    ("📡", "毎朝6時、AIニュース3本を自動収集・厳選", ACCENT_BLUE),
    ("📖", "4つのストーリーテリングフレームワークで「読みたくなるニュース」に変換", ACCENT_GREEN),
    ("📧", "Gmail (HTML) + LINE Notify でマルチチャネル配信", ACCENT_ORANGE),
    ("💡", "Markdownナレッジベースに蓄積。タグ検索・全文検索・月次サマリ対応", ACCENT_PURPLE),
    ("📊", "カウンターベースのリアクション機能でフィードバックループを実現", ACCENT_RED),
    ("✅", "219テスト全パス。品質保証済みの本番対応コード", ACCENT_TEAL),
]

for i, (icon, text, color) in enumerate(summary_items):
    y = Inches(2.0) + Inches(i * 0.8)

    # Icon circle
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(1.5), y, Inches(0.6), Inches(0.6))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_body_text(slide, text,
                  Inches(2.4), y + Inches(0.05),
                  Inches(9.0), Inches(0.5),
                  font_size=18, color=DARK_GRAY)

# Bottom tagline
add_body_text(slide, "AI News Collector Bot  —  AIの最前線を、毎朝あなたの手元に。",
              Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.5),
              font_size=20, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom accent line
bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(7.42), SLIDE_W, Inches(0.08))
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT_BLUE
bar2.line.fill.background()

add_page_number(slide, page)


# ── Save ────────────────────────────────────────────────────────
output_path = "/home/user/AI-News-Collection-Bot_v1/AI_News_Collector_Bot_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {page}")
